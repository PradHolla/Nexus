import os
import json
import re
import boto3
import subprocess
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Dict, Any, Tuple
import concurrent.futures
from src.job_status_store import append_job_log

bedrock_client = boto3.client('bedrock-runtime', region_name=os.getenv("AWS_REGION", "us-east-1"))
VISION_MODEL_ID = "moonshotai.kimi-k2.5"

CHARACTER_THRESHOLD = 300 

ADMIN_KEYWORDS = ["office hours", "grading policy", "zoom link", "late submission", "prerequisites", "textbooks", "ta information", "course work", "course overview"]

def convert_ppt_to_pptx(file_path: str, job_id: str = None) -> str:
    """Converts legacy .ppt to .pptx using LibreOffice headless."""
    if not file_path.lower().endswith(".ppt"):
        return file_path
    
    msg = f"Legacy .ppt detected. Converting {os.path.basename(file_path)} to .pptx via LibreOffice..."
    if job_id:
        append_job_log(job_id, msg)
    else:
        print(msg)
        
    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pptx", file_path, "--outdir", "/tmp/"],
            check=True,
            capture_output=True
        )
        base_name = os.path.basename(file_path)
        new_name = os.path.splitext(base_name)[0] + ".pptx"
        new_path = os.path.join("/tmp", new_name)
        
        if os.path.exists(new_path):
            return new_path
        else:
            raise FileNotFoundError("LibreOffice conversion failed to produce a file.")
    except Exception as e:
        msg = f"Failed to convert .ppt to .pptx: {e}"
        if job_id:
            append_job_log(job_id, msg)
        else:
            print(msg)
        raise e

def clean_json_response(content: str) -> dict:
    try:
        json_match = re.search(r'\{.*\}', content, re.DOTALL)
        if json_match:
            return json.loads(json_match.group(0))
        return {}
    except Exception as e:
        print(f"JSON Parsing Error: {e}")
        return {}

def sanitize_extracted_text(text: str) -> str:
    """Removes massive unbroken strings (like Base64 image leaks or hex dumps) that crash tokenizers."""
    cleaned = re.sub(r'\S{200,}', '[GARBAGE_DATA_REMOVED]', text)
    return cleaned.strip()

def analyze_diagram_with_kimi(image_bytes: bytes, image_ext: str, job_id: str = None) -> Dict[str, str]:
    prompt = """
    You are an expert academic tutor. Analyze this diagram, chart, or visual slide. 
    1. Describe the educational concept shown in high detail so it can be used for a quiz.
    2. Provide a 1-3 word topic tag for what this image represents.
    
    Return ONLY valid JSON. Format:
    {
        "description": "A detailed explanation...",
        "topic_tag": "Neural Networks"
    }
    """
    
    fmt = "png" if "png" in image_ext.lower() else "jpeg"
    
    try:
        response = bedrock_client.converse(
            modelId=VISION_MODEL_ID,
            messages=[{
                "role": "user",
                "content": [
                    {"image": {"format": fmt, "source": {"bytes": image_bytes}}},
                    {"text": prompt}
                ]
            }],
            inferenceConfig={"maxTokens": 512, "temperature": 0.1}
        )
        content = response['output']['message']['content'][0]['text']
        
        parsed_json = clean_json_response(content)
        return parsed_json if parsed_json else {"description": "Visual could not be parsed.", "topic_tag": "Diagram"}
        
    except Exception as e:
        msg = f"Vision fallback failed: {e}"
        if job_id:
            append_job_log(job_id, msg)
        else:
            print(msg)
        return {"description": "Vision API failed.", "topic_tag": "Diagram"}

def _extract_text_and_images(shape, image_blobs: List[Tuple[bytes, str]]) -> str:
    """Recursively extracts text from shapes, tables, and groups. Also collects images."""
    text_parts = []
    
    # 1. Direct Text
    if hasattr(shape, "text") and shape.text:
        text_parts.append(shape.text)
        
    # 2. Tables
    if shape.has_table:
        table_text = []
        for row in shape.table.rows:
            row_text = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame]
            table_text.append(" | ".join(row_text))
        text_parts.append("\n[Table Data]:\n" + "\n".join(table_text))
        
    # 3. Groups
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            text_parts.append(_extract_text_and_images(sub_shape, image_blobs))
            
    # 4. Images
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image_blobs.append((shape.image.blob, shape.image.ext))
        
    return "\n".join(text_parts)

def parse_ppt(file_path: str, filename: str, course_id: str, lecture_number: int, job_id: str = None) -> List[Dict[str, Any]]:
    # Handle legacy .ppt
    effective_path = file_path
    is_converted = False
    if file_path.lower().endswith(".ppt"):
        effective_path = convert_ppt_to_pptx(file_path, job_id)
        is_converted = True

    if not os.path.exists(effective_path):
        raise FileNotFoundError(f"Cannot find {effective_path}")

    msg = f"Processing Presentation: {filename}"
    if job_id:
        append_job_log(job_id, msg)
    else:
        print(msg)
    
    try:
        prs = Presentation(effective_path)
    except Exception as e:
        if is_converted and os.path.exists(effective_path):
            os.remove(effective_path)
        raise e

    chunks = []
    current_topic = "General Context"
    vision_tasks = [] 
    
    for slide_idx, slide in enumerate(prs.slides):
        page_num = slide_idx + 1
        
        image_blobs = []
        text_runs = []
        
        # 1. Extract from all shapes
        for shape in slide.shapes:
            text_runs.append(_extract_text_and_images(shape, image_blobs))
            
        # 2. Extract Speaker Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                text_runs.append(f"\n[Speaker Notes]:\n{notes}")
                
        clean_text = sanitize_extracted_text("\n".join(text_runs).strip())
        is_admin = any(kw in clean_text.lower() for kw in ADMIN_KEYWORDS)
        
        chunk_data = {
            "chunk_id": f"{course_id}_{filename}_s{page_num}",
            "file_name": filename,
            "course_id": course_id,
            "lecture_number": lecture_number,
            "page_number": page_num,
            "is_administrative": is_admin,
            "has_math": bool(re.search(r'[\$\\]', clean_text)),
            "used_vision": False
        }

        if is_admin:
            chunk_data["text"] = clean_text
            chunk_data["topic"] = "Course Administration"
            chunks.append(chunk_data)
            continue

        if slide.shapes.title and slide.shapes.title.text:
            current_topic = slide.shapes.title.text.strip()

        chunk_data["text"] = clean_text
        chunk_data["topic"] = current_topic

        # 3. Vision Trigger: Always pick largest image if ANY images found
        if image_blobs:
            largest_image, ext = max(image_blobs, key=lambda x: len(x[0]))
            chunk_data["used_vision"] = True
            
            vision_tasks.append({
                "chunk_ref": chunk_data, 
                "image_bytes": largest_image,
                "ext": ext,
                "page_num": page_num
            })
            
        chunks.append(chunk_data)

    # 4. Batch Process Vision
    if vision_tasks:
        msg = f"Concurrently processing {len(vision_tasks)} vision fallbacks via Kimi 2.5..."
        if job_id:
            append_job_log(job_id, msg)
        else:
            print(msg)
        
        def _fetch_vision(task):
            result = analyze_diagram_with_kimi(task["image_bytes"], task["ext"], job_id=job_id)
            return task, result

        with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
            futures = [executor.submit(_fetch_vision, task) for task in vision_tasks]
            
            for future in concurrent.futures.as_completed(futures):
                task, vision_result = future.result()
                chunk_ref = task["chunk_ref"]
                original_text = chunk_ref["text"]
                chunk_ref["text"] = f"{original_text}\n\n[Visual Description]: {vision_result.get('description', '')}"
                
                if vision_result.get("topic_tag") and vision_result.get("topic_tag") != "Diagram":
                    chunk_ref["topic"] = vision_result.get("topic_tag")
                    
        msg = "Vision processing complete."
        if job_id:
            append_job_log(job_id, msg)
        else:
            print(msg)

    # Cleanup temporary .pptx if it was converted
    if is_converted and os.path.exists(effective_path):
        try:
            os.remove(effective_path)
        except:
            pass

    return chunks
